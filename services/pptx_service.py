# PowerPoint processing service (REUSABLE)
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import logging
from pathlib import Path
from typing import List, Dict, Optional
from utils.profiler import profile
import subprocess
import platform
import shutil

logger = logging.getLogger(__name__)


class PPTXService:
    
    @staticmethod
    @profile
    def extract_images_from_slides(
        pptx_path: str,
        output_dir: str,
        max_slides: Optional[int] = None
    ) -> List[Dict]:
        """
        Extract images from PPTX slides (including groups)
        Returns list of: [{slide_idx, shape_id, file_path}, ...]
        """
        os.makedirs(output_dir, exist_ok=True)
        
        prs = Presentation(pptx_path)
        image_metadata = []
        
        slides_to_process = list(prs.slides)[:max_slides] if max_slides else prs.slides
        
        for slide_idx, slide in enumerate(slides_to_process):
            PPTXService._extract_from_shapes(
                shapes=slide.shapes,
                slide_idx=slide_idx,
                output_dir=output_dir,
                image_metadata=image_metadata
            )
        
        logger.info(f"Extracted {len(image_metadata)} images from {len(slides_to_process)} slides")
        return image_metadata
    
 
    @staticmethod
    def _extract_from_shapes(shapes, slide_idx: int, output_dir: str, image_metadata: List):
        """Recursively extract images from shapes (including groups)"""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    filename = f"slide{slide_idx}_shape{shape.shape_id}.png"
                    filepath = os.path.join(output_dir, filename)
                    
                    with open(filepath, 'wb') as f:
                        f.write(shape.image.blob)
                    
                    image_metadata.append({
                        'slide_idx': slide_idx,
                        'shape_id': shape.shape_id,
                        'file_path': filepath
                    })
                    
                except Exception as e:
                    logger.error(f"Error extracting image: {e}")
            
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recursively process groups
                PPTXService._extract_from_shapes(
                    shapes=shape.shapes,
                    slide_idx=slide_idx,
                    output_dir=output_dir,
                    image_metadata=image_metadata
                )
   
    @staticmethod
    @profile
    def replace_text_in_pptx(
        pptx_path: str,
        replacements: Dict[str, str],
        output_path: str
    ) -> str:
        """
        Replace text in PPTX and save as new file
        replacements: {old_text: new_text}
        """
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for old_text, new_text in replacements.items():
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
        
        prs.save(output_path)
        logger.info(f"Text replaced and saved: {output_path}")
        return output_path
    
   
    @staticmethod
    @profile
    def convert_slides_to_images(
        pptx_path: str,
        output_dir: str,
        max_slides: Optional[int] = None
    ) -> List[str]:
        """
        Convert PPTX slides to full-slide images (text + images combined)
        Direct PPTX → PNG conversion (no PDF intermediate step)
        Returns list of RELATIVE image file paths
        """
        import subprocess
        import platform
        import shutil
        import glob
        
        # Keep paths relative
        os.makedirs(output_dir, exist_ok=True)
        
        # Detect LibreOffice executable based on OS
        system = platform.system()
        
        if system == "Windows":
            soffice_cmd = shutil.which("soffice")
            if not soffice_cmd:
                common_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                ]
                for path in common_paths:
                    if os.path.exists(path):
                        soffice_cmd = path
                        break
            
            if not soffice_cmd:
                raise FileNotFoundError(
                    "LibreOffice not found. Please install LibreOffice and ensure "
                    "soffice.exe is in your PATH or installed in a standard location."
                )
        else:
            soffice_cmd = shutil.which("libreoffice") or shutil.which("soffice")
            
            if not soffice_cmd:
                raise FileNotFoundError(
                    "LibreOffice not found. Please install LibreOffice: "
                    "sudo apt-get install libreoffice (Ubuntu/Debian) or "
                    "brew install libreoffice (macOS)"
                )
        
        # Convert PPTX → PNG directly (150 DPI for preview speed)
        pptx_basename = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, f"{pptx_basename}.pdf")
        try:
            subprocess.run([
                soffice_cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                pptx_path
            ], check=True, capture_output=True, text=True)
            
            # Step 2: Convert PDF pages → PNG images using Ghostscript
            # Step 2: Convert PDF pages → PNG images using Ghostscript
            # Step 2: Convert PDF pages → PNG images using pdf2image
            from pdf2image import convert_from_path

            images = convert_from_path(
                pdf_path, 
                dpi=150,
                first_page=1,
                last_page=max_slides if max_slides else None
            )

            # Save each page as PNG
            image_paths = []
            for idx, image in enumerate(images):
                final_path = os.path.join(output_dir, f"slide_{idx}.png")
                image.save(final_path, 'PNG')
                image_paths.append(final_path)

            # Clean up PDF file
            if os.path.exists(pdf_path):
                os.remove(pdf_path)

            logger.info(f"Converted {len(image_paths)} slides to images (PPTX→PDF→PNG)")
            return image_paths

        except Exception as e:
            logger.error(f"Conversion failed: {e}")
            raise




    
    @staticmethod
    @profile
    def replace_images_in_pptx(
        pptx_path: str,
        image_metadata: List[Dict],
        output_path: str
    ):
        """
        Replace image blobs in PPTX with swapped versions
        image_metadata: [{slide_idx, shape_id, swapped_path}, ...]
        """
    
        
        prs = Presentation(pptx_path)
        replaced_count = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            replaced_count += PPTXService._replace_in_shapes(
                shapes=slide.shapes,
                slide_idx=slide_idx,
                slide=slide,
                image_metadata=image_metadata
            )
        
        prs.save(output_path)
        logger.info(f"Replaced {replaced_count} images in PPTX: {output_path}")

    @staticmethod
    def _replace_in_shapes(shapes, slide_idx: int, slide, image_metadata: List[Dict]) -> int:
        """Recursively replace images in shapes (including groups)"""
        from pptx.parts.image import Image as PptxImage
        
        count = 0
        
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Find matching metadata
                matching = next(
                    (img for img in image_metadata 
                    if img['slide_idx'] == slide_idx and img['shape_id'] == shape.shape_id),
                    None
                )
                
                if matching and os.path.exists(matching['swapped_path']):
                    try:
                        # Get the image part
                        slide_part = slide.part
                        rId = shape._element.blip_rId
                        image_part = slide_part.related_part(rId)
                        
                        # Replace blob with swapped image
                        new_pptx_img = PptxImage.from_file(matching['swapped_path'])
                        image_part.blob = new_pptx_img._blob
                        
                        logger.debug(f"Replaced image: slide {slide_idx + 1}, shape {shape.shape_id}")
                        count += 1
                        
                    except Exception as e:
                        logger.error(f"Error replacing image: {e}")
            
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                count += PPTXService._replace_in_shapes(
                    shapes=shape.shapes,
                    slide_idx=slide_idx,
                    slide=slide,
                    image_metadata=image_metadata
                )
        
        return count



    @staticmethod
    @profile
    def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
        """Convert PPTX to PDF using LibreOffice"""
                
        os.makedirs(output_dir, exist_ok=True)
        
        system = platform.system()
        
        if system == "Windows":
            soffice_cmd = shutil.which("soffice")
            if not soffice_cmd:
                common_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                ]
                for path in common_paths:
                    if os.path.exists(path):
                        soffice_cmd = path
                        break
            if not soffice_cmd:
                raise FileNotFoundError("LibreOffice not found")
        else:
            soffice_cmd = shutil.which("libreoffice") or shutil.which("soffice")
            if not soffice_cmd:
                raise FileNotFoundError("LibreOffice not found")
        
        try:
            subprocess.run([
                soffice_cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                pptx_path
            ], check=True, capture_output=True, text=True)
            
            pptx_basename = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(output_dir, f"{pptx_basename}.pdf")
            
            if not os.path.exists(pdf_path):
                raise FileNotFoundError(f"PDF not generated: {pdf_path}")
            
            logger.info(f"PDF generated: {pdf_path}")
            return pdf_path
            
        except subprocess.CalledProcessError as e:
            logger.error(f"PDF conversion failed: {e.stderr}")
            raise
